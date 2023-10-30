import re
from pptx import Presentation
from pptx.util import Pt
from buildingdata import data_path
from buildingdata.logger import logger
from .utils import text_placeholder_map, row_map
from buildingdata.maps.maps import query_image
from pathlib import Path


class PresentationChurn:
    def __init__(self, template_ppt, df) -> None:
        """
        Presentation constructor
        """
        self.prs = Presentation(template_ppt)
        self.df = df
        self.df_dict = df.to_dict()
        self.font_size = 12

    def analyze_input_ppt(self):
        """
        Take the input file and analyze the structure.
        The output file contains marked up information to make it easier for
        generating future powerpoint templates.
        """
        index = 1
        slide_layout = self.prs.slide_layouts[index]
        slide = self.prs.slides.add_slide(slide_layout)
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = f"Title for Layout {index}"
        except AttributeError:
            print(f"No Title for Layout {index}")
        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is a special placeholder
                try:
                    if "Title" not in shape.text:
                        shape.text = (
                            f"Placeholder index:{phf.idx} type:{shape.name}"
                        )
                except AttributeError:
                    print(f"{phf.type} has no text attribute")

                print(f"{phf.idx} {shape.name}")

    def generate_slides_v0(self, addresses: list):
        """
        Generate slide function no street view.
        """
        if addresses is None:
            addresses = list(self.df[("ID", "addresse")])
        index = 1
        slide_layout = self.prs.slide_layouts[index]
        for i, address in enumerate(addresses):
            if i % 100 == 0:
                logger.info(
                    f"{i} addresses generated {self.df.shape[0] - i} remaining"
                )
            slide = self.prs.slides.add_slide(slide_layout)
            for shape in slide.placeholders:
                if "Title" in shape.name:
                    shape.text = f"{address}"
                elif "Picture Placeholder" in shape.name:
                    if shape.name[-1] == "3":
                        image_path = self.get_sat_image(i)
                        shape.insert_picture(image_path)
                elif "Text Placeholder" in shape.name:
                    number = (
                        shape.name[-2:]
                        if shape.name[-2].isnumeric()
                        else shape.name[-1]
                    )
                    shape.text = text_placeholder_map[number]
                    shape.text_frame.paragraphs[0].font.size = Pt(16)
                elif "Table Placeholder" in shape.name:
                    number = (
                        shape.name[-2:]
                        if shape.name[-2].isnumeric()
                        else shape.name[-1]
                    )
                    shape = shape.insert_table(rows=row_map[number], cols=2)
                    table = shape.table
                    if number == "8":
                        table = self.get_data("0.Localisation", i, table)
                    if number == "9":
                        table = self.get_data("2.Donnees_general", i, table)
                    if number == "10":
                        table = self.get_data("3.Donnees_systeme", i, table)
                    if number == "11":
                        table = self.get_data("4.Donnees_parois", i, table)
                    if number == "12":
                        table = self.get_data(
                            "5.Donnees_menuiseries", i, table
                        )
                    if number == "14":
                        table = self.get_data("1.Occupants", i, table)
                    table.first_row = False

    def generate_slides_v1(self, addresses: list):
        """
        Generate slide function with street view.
        """
        if addresses is None:
            addresses = list(self.df[("ID", "addresse")])
        index = 1
        slide_layout = self.prs.slide_layouts[index]
        for i, address in enumerate(addresses):
            if i % 100 == 0:
                logger.info(
                    f"{i} addresses generated {self.df.shape[0] - i} remaining"
                )
            slide = self.prs.slides.add_slide(slide_layout)
            for shape in slide.placeholders:
                number = (
                    shape.name[-2:]
                    if shape.name[-2].isnumeric()
                    else shape.name[-1]
                )
                if "Title" in shape.name:
                    shape.text = f"{address}"
                elif "Picture Placeholder" in shape.name:
                    if number == "3":
                        image_path = self.get_sat_image(i)
                    if number == "15":
                        image_path = self.get_sv_image(i)
                    shape.insert_picture(image_path)
                elif "Text Placeholder" in shape.name:
                    shape.text = text_placeholder_map[number]
                    shape.text_frame.paragraphs[0].font.size = Pt(16)
                elif "Table Placeholder" in shape.name:
                    number = (
                        shape.name[-2:]
                        if shape.name[-2].isnumeric()
                        else shape.name[-1]
                    )
                    shape = shape.insert_table(rows=row_map[number], cols=2)
                    table = shape.table
                    if number == "8":
                        table = self.get_data("0.Localisation", i, table)
                    if number == "9":
                        table = self.get_data("2.Donnees_general", i, table)
                    if number == "10":
                        table = self.get_data("3.Donnees_systeme", i, table)
                    if number == "11":
                        table = self.get_data("4.Donnees_parois", i, table)
                    if number == "12":
                        table = self.get_data(
                            "5.Donnees_menuiseries", i, table
                        )
                    if number == "14":
                        table = self.get_data("1.Occupants", i, table)
                    table.first_row = False

    def get_sv_image(self, i: int) -> str:
        """
        Getter function for street view images
        """
        if self.df[("2.Donnees_general", "streetview status")][i] is False:
            logger.info("No street view image available.")
            return data_path["images"] + "/not_available.jpg"

        addr = self.df[("ID", "addresse")][i]
        file = Path(data_path["street_view"] + f"/{addr}/gsv_0.jpg")
        try:
            _ = file.resolve(strict=True)
        except FileNotFoundError:
            query_image(addr)
            return self.get_sv_image(i)
        else:
            res = data_path["street_view"] + f"/{addr}/gsv_0.jpg"
            return res

    def get_sat_image(self, i: int) -> str:
        """
        Getter function for street view images
        """
        file = data_path["images"]
        file += f"/solar/{self.df[('ID', 'addresse')][i]}.jpg"
        file = Path(file)
        try:
            _ = file.resolve(strict=True)
        except FileNotFoundError:
            return data_path["images"] + "/not_available.jpg"
        else:
            res = data_path["images"]
            res += f"/solar/{self.df[('ID', 'addresse')][i]}.jpg"
            return res

    def get_data(self, name: str, i: int, table):
        if name == "0.Localisation":
            keys = self.get_localisation_keys(name)
        else:
            keys = list(
                k
                for k in self.df_dict.keys()
                if name in k and "flag DPE 2012" not in k
            )

        for j, key in enumerate(keys):
            table.cell(j, 0).text = key[1]
            table.cell(j, 1).text = str(self.df_dict[key][i])

            table.cell(j, 0).text_frame.paragraphs[0].font.size = Pt(
                self.font_size
            )
            table.cell(j, 1).text_frame.paragraphs[0].font.size = Pt(
                self.font_size
            )
            table.rows[j].height = Pt(15)
            if "4.Donnees_parois" in key:
                table.columns[0].width = Pt(155)

        return table

    def treat_string(self, s: str) -> str:
        """
        Remove unecessary characters from string
        """
        # Remove square brackets []
        s = re.sub(r"\[|\]", "", s)

        # Remove double quotes ""
        s = re.sub(r'"', "", s)

        # Remove trailing .0 from numbers
        s = re.sub(r"(\d+)\.0\b", r"\1", s)

        return s

    def treat_vitrage(self, s: str) -> str:
        """
        Formats the treat vitrage data by remplacing boolean values by YES/NO.

        Parameters:
            s: str
                String representation of the boolean to be formated

        Returns:
            Oui | Non: str
                Formatted string
        """
        return "Oui" if float(s) else "Non"

    def get_localisation_keys(self, name: str) -> list:
        """
        Gets the columns of table 0.Localisation that will be shown on slide.
        USeful because there is no need to present postal address data.

        Parameters:
            name: str
                Name of the localisation table.

        Returns:
            keys: list
                List of data keys for table 0.Localisation.

        """
        keys = list(
            k
            for k in self.df_dict.keys()
            if name in k
            and k[1] in ("altitude (m)", "distance bati. Historique (m)")
        )
        return keys
